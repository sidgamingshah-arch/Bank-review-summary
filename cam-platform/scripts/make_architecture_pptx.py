"""Generate an editable PowerPoint of the CAM Platform technical architecture.

    python scripts/make_architecture_pptx.py [output.pptx]

Native shapes/tables (fully editable in PowerPoint/Keynote/Google Slides), in the
product's indigo identity. Mirrors docs/architecture.md. Requires python-pptx.
"""
from __future__ import annotations

import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---- palette (the product's indigo) ---------------------------------------
INDIGO   = RGBColor(0x5A, 0x48, 0xE0)
INDIGO_D = RGBColor(0x4A, 0x39, 0xCF)
TINT     = RGBColor(0xEF, 0xEC, 0xFD)
INK      = RGBColor(0x1B, 0x17, 0x40)
INK_SOFT = RGBColor(0x42, 0x3D, 0x63)
MUTED    = RGBColor(0x73, 0x6E, 0x93)
GROUND   = RGBColor(0xF5, 0xF5, 0xFB)
SURFACE  = RGBColor(0xFF, 0xFF, 0xFF)
BORDER   = RGBColor(0xD6, 0xD2, 0xEA)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GREEN    = RGBColor(0x0F, 0x8A, 0x52)
AMBER    = RGBColor(0xA9, 0x69, 0x0A)
BLUE     = RGBColor(0x2F, 0x5F, 0xE0)
RED      = RGBColor(0xD3, 0x3B, 0x40)
GREENT   = RGBColor(0xE6, 0xF6, 0xEE)
AMBERT   = RGBColor(0xFB, 0xF0, 0xDA)
BLUET    = RGBColor(0xE7, 0xED, 0xFD)
REDT     = RGBColor(0xFC, 0xEB, 0xEC)

FONT = "Segoe UI"
MONO = "Consolas"
EMU = 914400
W, H = 13.333, 7.5

prs = Presentation()
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)
BLANK = prs.slide_layouts[6]


def _fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def _line(shape, color, width=1.0, dash=None):
    shape.line.color.rgb = color
    shape.line.width = Pt(width)
    if dash is not None:
        # set dash style via XML (python-pptx has no enum for this)
        from pptx.oxml.ns import qn
        ln = shape.line._get_or_add_ln()
        d = ln.find(qn("a:prstDash"))
        if d is None:
            d = ln.makeelement(qn("a:prstDash"), {})
            ln.append(d)
        d.set("val", dash)


def _para(tf, text, size=13, color=INK, bold=False, name=FONT, align=PP_ALIGN.LEFT,
          space_after=2, space_before=0, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    # support inline "mono spans" wrapped in backticks
    parts = text.split("`")
    for i, seg in enumerate(parts):
        if not seg:
            continue
        r = p.add_run()
        r.text = seg
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = MONO if i % 2 else name
    return p


def box(slide, x, y, w, h, fill=SURFACE, line=BORDER, line_w=1.0, radius=True, dash=None):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    _fill(shp, fill)
    _line(shp, line, line_w, dash)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, Pt(6))
    return shp


def textbox(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    return tb.text_frame


def arrow(slide, x1, y1, x2, y2, color=INDIGO, width=1.75, dash=None):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1),
                                   Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(width)
    if dash:
        _line(c, color, width, dash)
    # arrowhead
    from pptx.oxml.ns import qn
    ln = c.line._get_or_add_ln()
    tail = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
    ln.append(tail)
    return c


def slide_base(title, eyebrow, subtitle=None):
    s = prs.slides.add_slide(BLANK)
    bg = box(s, -0.1, -0.1, W + 0.2, H + 0.2, fill=GROUND, line=GROUND, radius=False)
    bg.shadow.inherit = False
    # accent rule + eyebrow + title
    rule = box(s, 0.55, 0.5, 0.09, 0.62, fill=INDIGO, line=INDIGO, radius=False)
    rule.shadow.inherit = False
    tf = textbox(s, 0.8, 0.42, W - 1.4, 1.0)
    _para(tf, eyebrow.upper(), size=11, color=INDIGO, bold=True, name=MONO, first=True)
    _para(tf, title, size=27, color=INK, bold=True, space_before=1)
    if subtitle:
        stf = textbox(s, 0.8, 1.28, W - 1.6, 0.5)
        _para(stf, subtitle, size=13, color=MUTED, first=True)
    return s


def footer(s):
    tf = textbox(s, 0.8, H - 0.5, W - 1.6, 0.35)
    _para(tf, "CAM Platform · Technical Architecture", size=9, color=MUTED, name=MONO, first=True)


def tile(slide, x, y, w, h, head, body, accent=INDIGO, head_size=13, body_size=11):
    shp = box(slide, x, y, w, h)
    tf = shp.text_frame
    tf.vertical_anchor = MSO_ANCHOR.TOP
    _para(tf, head, size=head_size, color=accent, bold=True, first=True)
    for ln in (body if isinstance(body, list) else [body]):
        _para(tf, ln, size=body_size, color=INK_SOFT, space_before=1)
    return shp


def add_table(slide, rows, x, y, w, col_w, header=True, fs=11):
    nrows, ncols = len(rows), len(rows[0])
    gt = slide.shapes.add_table(nrows, ncols, Inches(x), Inches(y), Inches(w),
                                Inches(0.4 * nrows)).table
    gt.first_row = header
    gt.horz_banding = False
    for j, cw in enumerate(col_w):
        gt.columns[j].width = Inches(cw)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = gt.cell(i, j)
            cell.margin_left = Pt(7); cell.margin_right = Pt(7)
            cell.margin_top = Pt(3); cell.margin_bottom = Pt(3)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = INDIGO if (header and i == 0) else (
                SURFACE if i % 2 else GROUND)
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            parts = str(val).split("`")
            for k, seg in enumerate(parts):
                if not seg and len(parts) > 1:
                    continue
                r = p.add_run(); r.text = seg
                r.font.size = Pt(fs)
                r.font.name = MONO if k % 2 else FONT
                r.font.bold = bool(header and i == 0)
                r.font.color.rgb = WHITE if (header and i == 0) else INK_SOFT
    return gt


# ============================================================ slides ========

# 1 — title
s = prs.slides.add_slide(BLANK)
box(s, -0.1, -0.1, W + 0.2, H + 0.2, fill=INK, line=INK, radius=False).shadow.inherit = False
box(s, -0.1, H - 1.9, W + 0.2, 2.0, fill=INDIGO_D, line=INDIGO_D, radius=False).shadow.inherit = False
tf = textbox(s, 0.9, 2.0, W - 1.8, 3.0)
_para(tf, "TECHNICAL ARCHITECTURE", size=14, color=RGBColor(0xB9, 0xAF, 0xFF), bold=True, name=MONO, first=True)
_para(tf, "CAM Platform", size=52, color=WHITE, bold=True, space_before=6)
_para(tf, "AI-Assisted Credit Assessment Memo generation — how the product works, end to end.",
      size=18, color=RGBColor(0xCF, 0xCB, 0xF0), space_before=10)
tf2 = textbox(s, 0.9, H - 1.5, W - 1.8, 1.0)
_para(tf2, "9 services   ·   5 data planes   ·   4-agent pipeline per section   ·   5 versioned master types",
      size=13, color=WHITE, name=MONO, first=True)
_para(tf2, "Companion to docs/architecture.md · contracts.md · traceability.md", size=11,
      color=RGBColor(0xD9, 0xD5, 0xF5), space_before=6)

# 2 — at a glance
s = slide_base("How it fits together", "Mental model")
tiles = [
    ("Config > code", "Prompts, KPIs, templates and doctypes are business-authored masters — the bank evolves the product without a vendor release."),
    ("One door", "Every cross-component call traverses the gateway (the APIM stand-in): authN, throttling, correlation, access logs."),
    ("Closed model plane", "Only service identities reach `/api/genai`; the SPA can never reach a model endpoint (NFR-10)."),
    ("Reproducible by snapshot", "A run freezes master versions, preferences, the resolved bundle and external context — replay is deterministic."),
    ("Grounded, not fabricated", "Sections draft only from attributed source facts; a deterministic backstop flags untraceable numbers/dates."),
    ("Auditable end to end", "Hash-chained audit events + one correlation id span upload → generation → edit → export."),
]
cw, ch, gx, gy = 3.9, 1.32, 0.25, 0.25
x0, y0 = 0.8, 1.7
for i, (h, b) in enumerate(tiles):
    r, c = divmod(i, 3)
    tile(s, x0 + c * (cw + gx), y0 + r * (ch + gy), cw, ch, h, b, head_size=13, body_size=10.5)
footer(s)

# 3 — topology
s = slide_base("Topology — 9 services behind one gateway", "Component view",
               "Dotted arrows are out-of-gateway egress. Ports shown are the local-dev layout.")
spa = box(s, 0.8, 1.8, 2.0, 0.65, fill=SURFACE, line=BORDER)
_para(spa.text_frame, "React SPA", size=12, color=INK, bold=True, align=PP_ALIGN.CENTER, first=True)
gw = box(s, 3.4, 1.75, 3.0, 0.75, fill=INDIGO, line=INDIGO_D)
_para(gw.text_frame, "Gateway :8080", size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER, first=True)
_para(gw.text_frame, "authN · routing · serves SPA", size=9, color=RGBColor(0xE5, 0xE1, 0xFF), align=PP_ALIGN.CENTER)
arrow(s, 2.8, 2.12, 3.4, 2.12)
svcs = [("auth", ":8101"), ("master-config", ":8102"), ("document", ":8103"), ("tagging", ":8104"),
        ("orchestration", ":8105"), ("genai", ":8106"), ("output", ":8107"), ("audit", ":8108")]
sw, sh, sgx, sgy = 2.0, 0.62, 0.18, 0.22
sx0, sy0 = 0.8, 3.05                     # 4-col block spans 0.8 → 9.34
for i, (n, p) in enumerate(svcs):
    r, c = divmod(i, 4)
    b = box(s, sx0 + c * (sw + sgx), sy0 + r * (sh + sgy), sw, sh, fill=SURFACE, line=BORDER)
    accent = INDIGO if n in ("genai", "orchestration") else INK
    _para(b.text_frame, n, size=10.5, color=accent, bold=True, align=PP_ALIGN.CENTER, first=True)
    _para(b.text_frame, p, size=8.5, color=MUTED, name=MONO, align=PP_ALIGN.CENTER)
arrow(s, 4.9, 2.5, 4.9, 3.05)            # gateway -> service block
# side integrations (right column, clear of the 9.34 grid edge)
llm = box(s, 9.9, 3.05, 2.9, 0.62, fill=TINT, line=INDIGO)
_para(llm.text_frame, "Model endpoint", size=11, color=INDIGO_D, bold=True, align=PP_ALIGN.CENTER, first=True)
conn = box(s, 9.9, 3.89, 2.9, 0.62, fill=AMBERT, line=AMBER)
_para(conn.text_frame, "Connectors (3rd-party)", size=11, color=AMBER, bold=True, align=PP_ALIGN.CENTER, first=True)
store = box(s, 9.9, 4.73, 2.9, 0.62, fill=GROUND, line=BORDER)
_para(store.text_frame, "Blob + PostgreSQL", size=11, color=INK_SOFT, bold=True, align=PP_ALIGN.CENTER, first=True)
arrow(s, 9.34, 3.36, 9.9, 3.36)                          # genai -> model endpoint
arrow(s, 9.34, 4.2, 9.9, 4.2, color=AMBER, dash="dash")  # orchestration -> connectors (direct)
arrow(s, 9.34, 5.04, 9.9, 5.04, color=MUTED)             # document -> storage
tf = textbox(s, 0.8, 5.75, 11.9, 0.7)
_para(tf, "Same image per service in containers; the gateway also serves the built SPA, so the whole app is one origin.",
      size=11, color=MUTED, first=True)
footer(s)

# 4 — five planes
s = slide_base("The five data planes", "Security backbone",
               "Each arrow in the topology belongs to one plane, with its own trust and egress rules.")
rows = [["Plane", "Initiator", "Path", "Auth carried", "Egress"],
        ["Control", "SPA → services", "via gateway", "end-user JWT", "internal"],
        ["Model", "services → genai", "via gateway", "service token ONLY", "genai → LLM"],
        ["Retrieval", "document svc", "to embed / search", "backend key (env)", "embeddings / index"],
        ["Connector", "orchestration", "direct, out-of-gateway", "connector key ONLY", "third-party feeds"],
        ["Storage", "document svc", "adapter", "conn string (env)", "blob store"]]
add_table(s, rows, 0.8, 1.9, 11.7, [1.7, 2.5, 2.7, 2.7, 2.1], fs=11.5)
tf = textbox(s, 0.8, 5.3, 11.7, 0.7)
_para(tf, "The connector plane is deliberately credential-isolated: a third-party vendor never receives the internal service token.",
      size=12, color=INK_SOFT, first=True)
footer(s)

# 5 — masters
s = slide_base("Business logic lives in versioned masters", "Configuration over code")
mt = [("Templates", "Ordered sections; conditional, fixed-format, depends_on, uses_external_context"),
      ("Section prompts", "Per-section instructions · `{{placeholders}}` · `{{industry_kpis}}`"),
      ("Standing rules", "House style + guardrails prepended to every prompt"),
      ("Document types", "Doctypes for classification + section mapping"),
      ("KPI sets", "Per-industry KPI definitions rendered into the prompt")]
for i, (h, b) in enumerate(mt):
    tile(s, 0.8, 1.75 + i * 0.82, 7.6, 0.72, h, b, head_size=13, body_size=11)
mc = box(s, 8.7, 1.75, 3.8, 4.0, fill=TINT, line=INDIGO)
tfm = mc.text_frame; tfm.vertical_anchor = MSO_ANCHOR.TOP
_para(tfm, "Maker-checker", size=14, color=INDIGO_D, bold=True, first=True)
_para(tfm, "Every master is versioned. A maker drafts a new version; a different checker approves it (self-approval rejected).",
      size=11.5, color=INK_SOFT, space_before=6)
_para(tfm, "Only an approved/published version can be resolved by a run — otherwise the run refuses with `not_published`.",
      size=11.5, color=INK_SOFT, space_before=8)
footer(s)

# 6 — document lifecycle
s = slide_base("Document lifecycle", "Intake before grounding",
               "Every document — uploaded, pulled, or dropped into chat — passes the same intake.")
steps = [("VAF intake", "validate → AV scan → quarantine with a visible reason", AMBER, AMBERT),
         ("Extract", "text extract to blob storage; DB holds only metadata + hashes", BLUE, BLUET),
         ("Tag", "AI-first classify (keyword fallback); the method is audited", INDIGO, TINT),
         ("Index", "if RAG on: chunk → embed → vector store / Azure Search", GREEN, GREENT)]
bw, bx0, by = 2.85, 0.8, 2.4
for i, (h, b, ac, bg) in enumerate(steps):
    x = bx0 + i * (bw + 0.25)
    bxs = box(s, x, by, bw, 1.5, fill=bg, line=ac)
    tfb = bxs.text_frame; tfb.vertical_anchor = MSO_ANCHOR.TOP
    _para(tfb, h, size=14, color=ac, bold=True, first=True)
    _para(tfb, b, size=11, color=INK_SOFT, space_before=4)
    if i < 3:
        arrow(s, x + bw, by + 0.75, x + bw + 0.25, by + 0.75, color=MUTED)
tf = textbox(s, 0.8, 4.4, 11.7, 0.6)
_para(tf, "Quarantined content is never stored or used as grounding.", size=12, color=INK_SOFT, first=True)
footer(s)

# 7 — generation flow
s = slide_base("The generation flow", "Resolve → deliver",
               "One job per section; a worker pool drives a four-agent pipeline through the GenAI gateway.")
flow = ["Resolve", "Snapshot", "Queue", "Admit", "Pipeline", "Check", "Deliver"]
fw, fx0, fy = 1.55, 0.8, 2.05
for i, step in enumerate(flow):
    x = fx0 + i * (fw + 0.13)
    ac = INDIGO if step in ("Pipeline",) else INK
    fillc = TINT if step == "Pipeline" else SURFACE
    b = box(s, x, fy, fw, 0.7, fill=fillc, line=INDIGO if step == "Pipeline" else BORDER)
    _para(b.text_frame, f"{i+1}. {step}", size=11.5, color=ac, bold=True, align=PP_ALIGN.CENTER, first=True)
    if i < len(flow) - 1:
        arrow(s, x + fw, fy + 0.35, x + fw + 0.13, fy + 0.35, color=MUTED, width=1.4)
agents = [("Extraction", "source-attributed facts from only the mapped docs"),
          ("Summarisation", "drafts through the layered prompt + KPIs + style"),
          ("Materiality check", "coverage verdict → bounded revision loop"),
          ("Consistency check", "vs facts + other sections → revision loop")]
for i, (h, b) in enumerate(agents):
    r, c = divmod(i, 2)
    tile(s, 0.8 + c * 5.95, 3.25 + r * 1.15, 5.7, 1.02, h, b, head_size=12.5, body_size=11)
tf = textbox(s, 0.8, 5.7, 11.7, 0.7)
_para(tf, "Every agent call is recorded in the section's `agent_trace`. Failures stay section-local: retry or regenerate one section.",
      size=11, color=MUTED, first=True)
footer(s)

# 8 — connectors (1)
s = slide_base("Connectors — external intelligence", "Deep dive · 1 of 2",
               "Client feeds (negative-news, web/market context) become grounding — the model never calls out itself.")
seq = [("Gate", "section opts in\n(uses_external_context)\nAND admin toggle on", INDIGO, TINT),
       ("Fetch once", "per run, keyed by\nborrower + industry", INDIGO, SURFACE),
       ("Call vendor", "direct, out-of-gateway\nX-Connector-Key ONLY\n(or labelled MOCK)", AMBER, AMBERT),
       ("Snapshot", "into run.\nconnector_context", INDIGO, SURFACE),
       ("Ground + disclose", "opted-in sections read it;\nshown in gap trailer", GREEN, GREENT)]
bw = 2.25
for i, (h, b, ac, bg) in enumerate(seq):
    x = 0.8 + i * (bw + 0.2)
    bxs = box(s, x, 2.35, bw, 1.7, fill=bg, line=ac)
    tfb = bxs.text_frame; tfb.vertical_anchor = MSO_ANCHOR.TOP
    _para(tfb, h, size=12.5, color=ac, bold=True, first=True)
    for ln in b.split("\n"):
        _para(tfb, ln, size=10, color=INK_SOFT, space_after=0)
    if i < len(seq) - 1:
        arrow(s, x + bw, 3.2, x + bw + 0.2, 3.2, color=MUTED, width=1.4)
note = box(s, 0.8, 4.45, 11.7, 1.35, fill=AMBERT, line=AMBER)
tfn = note.text_frame; tfn.vertical_anchor = MSO_ANCHOR.TOP
_para(tfn, "Isolated egress (why it's out-of-gateway)", size=13, color=AMBER, bold=True, first=True)
_para(tfn, "A connector is a third-party host outside the bank gateway, so orchestration calls it directly with only its own connector key — never the internal service token. Handing a vendor a valid platform credential is precisely what this design avoids (NFR-06).",
      size=11.5, color=INK_SOFT, space_before=4)
footer(s)

# 9 — connectors (2)
s = slide_base("Connectors — properties & configuration", "Deep dive · 2 of 2")
props = [("Double-gated", "per-section opt-in AND admin master toggle (default off)", INDIGO),
         ("Fail-open", "any error/timeout/non-200 → empty; the run never fails on it", GREEN),
         ("Mock mode", "enabled but no URL → labelled MOCK item (offline demo/test)", BLUE),
         ("Injection-hardened", "labels stripped, count-capped, wrapped as inert document", RED)]
for i, (h, b, ac) in enumerate(props):
    r, c = divmod(i, 2)
    tile(s, 0.8 + c * 5.95, 1.8 + r * 1.15, 5.7, 1.02, h, b, accent=ac, head_size=12.5, body_size=11)
rows = [["Setting", "Meaning", "Default"],
        ["connectors_*_enabled", "admin toggles (news / search)", "off"],
        ["CAM_CONNECTOR_*_URL", "vendor endpoints; empty → mock", '""'],
        ["CAM_CONNECTOR_API_KEY_ENV", "name of the env var holding the key", "CAM_CONNECTOR_API_KEY"],
        ["CAM_CONNECTOR_TIMEOUT_SECONDS", "per-call timeout", "8"],
        ["CAM_CONNECTOR_MAX_ITEMS", "items kept per connector", "5"]]
add_table(s, rows, 0.8, 4.2, 11.7, [4.0, 5.2, 2.5], fs=11)
tf = textbox(s, 0.8, H - 0.92, 11.7, 0.4)
_para(tf, "Vendor contract:  POST {borrower, industry, max_items} → {items:[{title, source, date, text|summary}]}",
      size=11, color=INDIGO_D, name=MONO, first=True)

# 10 — RAG
s = slide_base("Retrieval / RAG", "Large documents",
               "For big files, grounding is retrieved rather than whole-document.")
tile(s, 0.8, 1.9, 5.7, 1.5, "Three modes (rag_mode)",
     ["off — whole-document grounding", "keyword — BM25 / lexical", "embedding — vector similarity"], body_size=12)
tile(s, 6.75, 1.9, 5.7, 1.5, "Backends (CAM_RETRIEVAL_BACKEND)",
     ["local — in-DB vector / lexical", "azure_search — Azure AI Search", "one /retrieve contract either way"], body_size=12)
tile(s, 0.8, 3.65, 5.7, 1.5, "Embedding egress",
     ["centralised at genai /embed (service-token only)", "OpenAI-compatible · Azure · local hash"], body_size=12)
tile(s, 6.75, 3.65, 5.7, 1.5, "Provenance & fallback",
     ["top-K passages recorded per section", "unreachable / unindexed → whole-doc fallback (never a hard failure)"], body_size=12)
footer(s)

# 11 — genai gateway
s = slide_base("GenAI gateway & providers", "The only component that talks to a model",
               "Service identities only — the gateway also blocks end-user tokens at the edge.")
gg = [("Assembly", "layered prompt (house → standing rules → template → section); grounding wrapped in inert document blocks"),
      ("Providers", "mock · anthropic · openai (any OpenAI-compatible) · azure — swappable at runtime; keys by env-var name only"),
      ("Task endpoints", "generate · extract · materiality · consistency · reconcile · classify · embed · edit (+ token accounting)"),
      ("Deterministic backstop", "numeric/date traceability flags that no prompt can disable")]
for i, (h, b) in enumerate(gg):
    tile(s, 0.8, 1.9 + i * 1.03, 11.7, 0.92, h, b, head_size=13, body_size=11.5)
footer(s)

# 12 — queue & concurrency
s = slide_base("Queue & concurrency", "Two independent levers")
tile(s, 0.8, 2.0, 5.7, 2.1, "Section concurrency",
     ["`worker_concurrency` — how many sections draft in parallel",
      "clamped to the worker pool spawned at deploy (`CAM_WORKER_POOL_SIZE`)",
      "applies at runtime, no restart"], body_size=12)
tile(s, 6.75, 2.0, 5.7, 2.1, "Run concurrency",
     ["`max_concurrent_runs` (default 4)",
      "a burst is ACCEPTED and queued — never 429'd",
      "admitted FIFO as slots free; per-user fairness cap"], body_size=12)
note = box(s, 0.8, 4.4, 11.7, 1.0, fill=GROUND, line=BORDER)
_para(note.text_frame, "Claims are serialised in-process by a lock and, across processes on PostgreSQL, by SELECT … FOR UPDATE SKIP LOCKED — so workers never double-process; stuck claims are reaped by lease timeout.",
      size=12, color=INK_SOFT, first=True)
footer(s)

# 13 — editing & notifications
s = slide_base("Human-in-the-loop editing & notifications", "The analyst stays in control")
tile(s, 0.8, 1.95, 11.7, 1.6, "Editing (output service owns the working copy)",
     ["Per-section versions (autosave, named versions, diffs, optimistic locking).",
      "Chat replies land as PENDING suggestions with a diff — accepted/rejected explicitly, never auto-applied.",
      "In-chat uploads pass the same VAF + tagging pipeline; finalisation blocked while a suggestion is pending; drafts watermarked."],
     head_size=13, body_size=11.5)
tile(s, 0.8, 3.75, 11.7, 1.6, "Notifications (on any terminal run state, best-effort)",
     ["In-app — a per-user bell that polls unread and deep-links to the run.",
      "Email — gated by a master toggle; sent on a background thread (a slow relay never stalls the worker).",
      "No SMTP host configured → the mailer logs instead of sending; password read from an env var by name (NFR-06)."],
     head_size=13, body_size=11.5)
footer(s)

# 14 — security
s = slide_base("Security model", "Fail-open at the edges, fail-closed at the core")
rows = [["Concern", "Mechanism"],
        ["Identity", "short-lived HS256 JWTs (dev IdP stub); production swaps one auth-adapter for the bank IdP"],
        ["Authorisation", "single role → capability matrix in every service; analysts own-scoped"],
        ["Model-plane closure", "/api/genai accepts service tokens only; gateway rejects user tokens at the edge"],
        ["Connector isolation", "third-party feeds get only their connector key, never the internal service token"],
        ["Secrets", "env/vault only; keys & SMTP password by env-var NAME, never on Settings or in logs"],
        ["Prompt injection", "all grounding sanitised + wrapped in inert document blocks; docs marked as data"],
        ["Malware", "VAF: validate → AV scan → quarantine with visible reason; never stored/used"],
        ["Tamper evidence", "audit events hash-chained sha256(prev + canonical(event)) + verify-chain"]]
add_table(s, rows, 0.8, 1.85, 11.7, [3.0, 8.7], fs=11)
footer(s)

# 15 — deployment
s = slide_base("Deployment — Linux & Windows", "One codebase, several targets",
               "Selected by environment, not by code. Gateway (and the built SPA) on :8080.")
dep = [("One-click dev", "Windows: start-windows.bat · Linux/macOS: ./start-linux.sh → venv + UI + seed + stack + browser", INDIGO, TINT),
       ("Linux service", "systemd unit (deploy/systemd/) supervising run_stack.py; restart-on-failure; SIGTERM-clean", GREEN, GREENT),
       ("Windows service", "boot Scheduled Task (deploy/windows/, built-in) or NSSM for a true SCM service", BLUE, BLUET),
       ("Containers", "docker compose — PostgreSQL + one container per service; scale orchestration horizontally", INDIGO, TINT),
       ("Cloud / bank", "Azure OpenAI · Azure AI Search · Azure Blob by env; behind the real APIM + bank IdP", AMBER, AMBERT)]
for i, (h, b, ac, bg) in enumerate(dep):
    bxs = box(s, 0.8, 1.85 + i * 0.85, 11.7, 0.76, fill=bg, line=ac)
    tfb = bxs.text_frame
    _para(tfb, h, size=12.5, color=ac, bold=True, first=True)
    _para(tfb, b, size=11, color=INK_SOFT, space_before=1)
footer(s)

# 16 — close
s = prs.slides.add_slide(BLANK)
box(s, -0.1, -0.1, W + 0.2, H + 0.2, fill=INK, line=INK, radius=False).shadow.inherit = False
tf = textbox(s, 0.9, 2.6, W - 1.8, 2.2)
_para(tf, "Grounded · reproducible · auditable", size=16, color=RGBColor(0xB9, 0xAF, 0xFF), bold=True, name=MONO, first=True)
_para(tf, "CAM Platform", size=40, color=WHITE, bold=True, space_before=6)
_para(tf, "Full detail: docs/architecture.md · docs/DEPLOYMENT.md · docs/contracts.md · docs/traceability.md",
      size=14, color=RGBColor(0xCF, 0xCB, 0xF0), space_before=12)


def main():
    out = sys.argv[1] if len(sys.argv) > 1 else "docs/CAM-Platform-Architecture.pptx"
    prs.save(out)
    print(f"wrote {out} ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
